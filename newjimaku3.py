from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
import pandas as pd

def create_presentation(text_list, image_paths, output_file="output_presentation_new.pptx"):
    prs = Presentation()

    # スライドサイズ設定
    prs.slide_width = Inches(13.3334646)
    prs.slide_height = Inches(7.5)

    font_size = 36  # フォントサイズ

    for i, (text, img_path) in enumerate(zip(text_list, image_paths)):
        # スライド生成スキップ条件
        if not text and not img_path:
            continue  # テキストも画像も指定されていない場合、スキップ

        slide = prs.slides.add_slide(prs.slide_layouts[7])

        # 真っ黒なスライドを作成する場合
        if img_path == "暗転":
            black_fill = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            black_fill.fill.solid()
            black_fill.fill.fore_color.rgb = RGBColor(0, 0, 0)
            continue  # 暗転スライドでは文字を設定しない

        # スライドに画像を挿入
        if img_path:
            try:
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            except FileNotFoundError:
                print(f"画像ファイルが見つかりません: {img_path}")

        # 白い背景のテキストエリア（透過度設定付き）
        if text:
            left = Inches(0)
            top = prs.slide_height * 3 / 4
            width = prs.slide_width
            height = prs.slide_height / 4

            white_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            white_box.fill.solid()
            white_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            white_box.fill.fore_color.transparency = 1  # 透過度を50%に設定

            # テキストボックスの挿入
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame

            text_frame.clear()
            text_frame.auto_size = MSO_AUTO_SIZE.NONE
            text_frame.word_wrap = True  # 自動改行

            # テキスト設定
            p = text_frame.add_paragraph()
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            #p.font.name = "MS Gothic"  # フォントをMSゴシックに設定

            # テキストの中央揃え
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # プレゼンテーションを保存
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

# エクセルファイルからデータを読み込む
file_path = "Book1.xlsx"
sheet_name = "Sheet1"

# 必要な列（A列:テキスト, B列:画像パス）をインデックスで読み込み
df = pd.read_excel(file_path, sheet_name=sheet_name, usecols=[0, 1], engine='openpyxl')

# テキストと画像パスのリストを作成
texts_raw = df.iloc[:, 0].fillna("").tolist()  # A列をテキストとして使用
images_raw = df.iloc[:, 1].fillna("").tolist()  # B列を画像パスとして使用

# A列の文字列を@で分割し、分割数に応じて画像リストを調整
texts = []
images = []

for text, image in zip(texts_raw, images_raw):
    if "@" in text:
        parts = text.split("@")  # @で分割
        texts.extend(parts)  # 分割したテキストを追加
        images.extend([image] * len(parts))  # 同じ画像パスを複製して対応付け
    else:
        texts.append(text)
        images.append(image)

# プレゼンテーションを作成
create_presentation(texts, images)
print("perfect")
